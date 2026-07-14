import argparse
import json
from pathlib import Path

from pptx import Presentation


def extract_pptx(path: str) -> dict:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    texts.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_number": i,
            "text": "\n".join(texts),
            "notes": notes,
        })

    return {"source": Path(path).name, "slides": slides}


def main():
    parser = argparse.ArgumentParser(description="Extract slide text from a .pptx file")
    parser.add_argument("pptx_path")
    parser.add_argument("-o", "--output", help="Write JSON to this path instead of stdout")
    args = parser.parse_args()

    result = extract_pptx(args.pptx_path)
    output = json.dumps(result, ensure_ascii=False, indent=2)

    if args.output:
        Path(args.output).write_text(output, encoding="utf-8")
    else:
        print(output)


if __name__ == "__main__":
    main()
