import pathlib
import sys

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parents[3]))  # repo root

from pptx import Presentation

from backend.stt.ppt_extractor import extract_pptx


def _make_sample_pptx(path):
    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "REST API 개요"
    slide1.placeholders[1].text = "REST API는 HTTP 기반의 아키텍처 스타일이다."
    slide1.notes_slide.notes_text_frame.text = "발표자 노트: 예시 위주로 설명"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "정리"
    slide2.placeholders[1].text = "핵심 개념 요약"

    prs.save(path)


def test_extract_pptx_returns_text_per_slide(tmp_path):
    pptx_path = tmp_path / "sample.pptx"
    _make_sample_pptx(pptx_path)

    result = extract_pptx(str(pptx_path))

    assert result["source"] == "sample.pptx"
    assert len(result["slides"]) == 2

    slide1 = result["slides"][0]
    assert slide1["slide_number"] == 1
    assert "REST API 개요" in slide1["text"]
    assert "HTTP 기반의 아키텍처 스타일" in slide1["text"]
    assert slide1["notes"] == "발표자 노트: 예시 위주로 설명"

    slide2 = result["slides"][1]
    assert slide2["slide_number"] == 2
    assert "정리" in slide2["text"]
    assert slide2["notes"] == ""


def test_extract_pptx_empty_deck(tmp_path):
    pptx_path = tmp_path / "empty.pptx"
    Presentation().save(pptx_path)

    result = extract_pptx(str(pptx_path))

    assert result["source"] == "empty.pptx"
    assert result["slides"] == []
